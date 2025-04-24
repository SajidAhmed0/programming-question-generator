import os
import glob
import PyPDF2
from pptx import Presentation
from tqdm import tqdm
from typing import List
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from pinecone import Pinecone
import openai

from dotenv import load_dotenv

load_dotenv()

# --- Configuration ---
openai.api_key = os.getenv("OPENAI_API_KEY")
pc = Pinecone(
        api_key=os.getenv("PINECONE_API_KEY")
    )

index_name = os.getenv("PINECONE_INDEX_NAME")

index = pc.Index(index_name)

BASE_DIR = "lecture_slides"
TEXT_DIR = "extracted_txt_files"
MODULE_TEXT_DIR = "merged_module_texts"

os.makedirs(TEXT_DIR, exist_ok=True)
os.makedirs(MODULE_TEXT_DIR, exist_ok=True)


# --- Extractor Class ---
class SlideExtractor:
    def __init__(self, base_dir: str, output_dir: str, merged_dir: str):
        self.base_dir = base_dir
        self.output_dir = output_dir
        self.merged_dir = merged_dir

    def extract_all(self):
        for module_dir in sorted(os.listdir(self.base_dir)):
            module_path = os.path.join(self.base_dir, module_dir)
            if not os.path.isdir(module_path):
                continue

            module_texts = []

            lecture_files = sorted(
                glob.glob(os.path.join(module_path, "*.pdf")) + glob.glob(os.path.join(module_path, "*.pptx")),
                key=lambda x: os.path.basename(x).lower()
            )

            for file_path in lecture_files:
                if file_path.endswith(".pdf"):
                    text = self._extract_pdf(file_path)
                else:
                    text = self._extract_pptx(file_path)

                if text:
                    self._save_individual_text(file_path, text)
                    module_texts.append(f"[{os.path.basename(file_path)}]\n{text}\n")

            # Save merged module-level file
            if module_texts:
                merged_path = os.path.join(self.merged_dir, f"{module_dir}.txt")
                with open(merged_path, "w", encoding="utf-8") as f:
                    f.write("\n".join(module_texts))

    def _extract_pdf(self, path: str) -> str:
        try:
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception as e:
            print(f"[PDF ERROR] {path} - {e}")
            return ""

    def _extract_pptx(self, path: str) -> str:
        try:
            prs = Presentation(path)
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        except Exception as e:
            print(f"[PPTX ERROR] {path} - {e}")
            return ""

    def _save_individual_text(self, original_path: str, content: str):
        rel_path = os.path.relpath(original_path, self.base_dir)
        txt_path = os.path.join(self.output_dir, rel_path.replace(".pdf", ".txt").replace(".pptx", ".txt"))
        os.makedirs(os.path.dirname(txt_path), exist_ok=True)
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write(content)


# --- Chunking & Embedding ---
class Vectorizer:
    def __init__(self, merged_dir: str, chunk_size: int = 1000, chunk_overlap: int = 100):
        self.merged_dir = merged_dir
        self.text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        self.embedder = OpenAIEmbeddings(model="text-embedding-3-small")

    def process(self):
        txt_files = sorted(glob.glob(os.path.join(self.merged_dir, "*.txt")))
        for path in tqdm(txt_files, desc="Vectorizing merged module files"):
            with open(path, "r", encoding="utf-8") as f:
                content = f.read()

            module_name = os.path.splitext(os.path.basename(path))[0]
            chunks = self.text_splitter.split_text(content)

            if chunks:
                for i, text in enumerate(chunks):
                    try:
                        vector = self.embedder.embed_query(text)  # using embed_query for single text
                        metadata = {
                            "module": module_name,
                            "source": f"{module_name}.txt",
                            "chunk_index": i,
                            "text": text,
                        }
                        index.upsert(
                            vectors=[(f"{module_name}-{i}", vector, metadata)],
                            namespace=module_name
                        )
                    except Exception as e:
                        print(f"[ERROR] Failed to process chunk {i} of {module_name}: {e}")



# --- Main Execution ---
if __name__ == "__main__":
    print("🔍 Extracting text from slides...")
    # extractor = SlideExtractor(BASE_DIR, TEXT_DIR)
    # extractor = SlideExtractor(BASE_DIR, TEXT_DIR, MODULE_TEXT_DIR)
    # extractor.extract_all()

    print("🧠 Vectorizing and storing into Pinecone...")
    # vectorizer = Vectorizer(MODULE_TEXT_DIR)
    # vectorizer.process()

    print("✅ All done!")
